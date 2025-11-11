import os
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


PROJECT_TITLE = "House Price Prediction"
AUTHOR = "KUSHAL SANGUR"
INTERNSHIP_ID = "ALPH08ML536"
ORG = "SSCASC COLLEGE TUMKUR"


def add_title_slide(prs: Presentation) -> None:
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = PROJECT_TITLE
    # Center title text
    if title_shape.has_text_frame:
        title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    subtitle = slide.placeholders[1]
    subtitle.text = f"{AUTHOR} | {INTERNSHIP_ID}\n{ORG}\n{datetime.now().strftime('%b %d, %Y')}"
    if subtitle.has_text_frame:
        for p in subtitle.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    text_frame = slide.shapes.placeholders[1].text_frame
    text_frame.clear()
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(20)


def add_picture_slide(prs: Presentation, title: str, image_path: str) -> None:
    slide_layout = prs.slide_layouts[5]  # Title Only
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = title
    if title_shape.has_text_frame:
        title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    if os.path.exists(image_path):
        # Dynamic centering and sizing
        slide_width = prs.slide_width
        margin = Inches(0.75)
        max_width = slide_width - 2 * margin
        # Cap max width to keep margins aesthetically pleasing
        hard_cap = Inches(9.5)
        img_width = max_width if max_width < hard_cap else hard_cap

        left = int((slide_width - img_width) / 2)
        top = Inches(1.6)

        pic = slide.shapes.add_picture(image_path, left, top, width=img_width)

        # Make image clickable (opens source file)
        try:
            pic.click_action.hyperlink.address = os.path.abspath(image_path)
        except Exception:
            pass

        # Caption centered under image
        caption_top = top + pic.height + Inches(0.2)
        caption_height = Inches(0.6)
        caption = slide.shapes.add_textbox(left, caption_top, img_width, caption_height)
        tf = caption.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "Click image to view full size"
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(80, 80, 80)
    else:
        body = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
        tf = body.text_frame
        p = tf.paragraphs[0]
        p.text = f"Image not found: {image_path}"
        p.font.color.rgb = RGBColor(200, 0, 0)


def build_presentation(output_path: str = "Project_Presentation.pptx") -> str:
    prs = Presentation()

    # 1) Title Slide
    add_title_slide(prs)

    # 2) Introduction & Problem Statement
    add_bullet_slide(
        prs,
        "Introduction & Problem Statement",
        [
            "Goal: Predict residential house prices using ML.",
            "Problem: Manual valuation is time-consuming and inconsistent.",
            "Impact: Faster, data-driven pricing for buyers, sellers, and agents.",
        ],
    )

    # 3) Literature Review / Existing System
    add_bullet_slide(
        prs,
        "Literature Review / Existing System",
        [
            "Hedonic pricing models (linear) widely used in real estate.",
            "Tree-based ensembles (RF/GB/XGBoost) outperform linear baselines.",
            "Existing tools lack feature engineering and robust evaluation.",
        ],
    )

    # 4) Proposed System / Methodology
    add_bullet_slide(
        prs,
        "Proposed System / Methodology",
        [
            "Data pipeline: generation → cleaning → feature engineering → scaling → split.",
            "Model zoo: Linear, Ridge, Lasso, Decision Tree, RF, GB, XGBoost, SVR, KNN.",
            "Model selection: Evaluate on RMSE, MAE, R², MAPE; tune best model.",
            "Deployment: Streamlit app for interactive predictions.",
        ],
    )

    # 5) Tools & Technologies Used
    add_bullet_slide(
        prs,
        "Tools & Technologies Used",
        [
            "Python, NumPy, Pandas",
            "Scikit-learn, XGBoost",
            "Matplotlib, Seaborn, Plotly",
            "Streamlit (UI)",
            "Joblib (model persistence)",
        ],
    )

    # 6) Implementation (step by step with visuals)
    add_bullet_slide(
        prs,
        "Implementation Steps",
        [
            "Data generation: house_data.csv via data_generator.py.",
            "Preprocessing: cleaning, engineered features, scaling, split.",
            "Training: multiple algorithms + hyperparameter tuning.",
            "Evaluation: metrics and visual comparisons.",
            "App: Streamlit interface for prediction and analysis.",
        ],
    )

    # Optional visuals during implementation
    if os.path.exists("data_analysis.png"):
        add_picture_slide(prs, "Data Analysis Visuals", "data_analysis.png")

    # 7) Results & Screenshots
    # Add model comparison and predictions vs actual plots
    add_picture_slide(prs, "Model Comparison (Results)", "model_comparison.png")
    add_picture_slide(prs, "Predictions vs Actual (Accuracy)", "predictions_vs_actual.png")

    # 8) Conclusion & Future Scope
    add_bullet_slide(
        prs,
        "Conclusion & Future Scope",
        [
            "Random Forest (tuned) achieved strong accuracy (see metrics).",
            "Feature engineering significantly improved performance.",
            "Future: real-time data, geographic features, forecasting, API & mobile app.",
        ],
    )

    # 9) References
    add_bullet_slide(
        prs,
        "References",
        [
            "Scikit-learn Documentation",
            "XGBoost Documentation",
            "Streamlit Documentation",
            "Project README for pipeline and metrics",
        ],
    )

    prs.save(output_path)
    return os.path.abspath(output_path)


if __name__ == "__main__":
    path = build_presentation()
    print(f"Presentation created: {path}")


